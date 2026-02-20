"""Report Agent — generates professional branded reports.

Supports PDF (WeasyPrint), PPTX (python-pptx), and Excel (openpyxl).
5 report types: executive, rssi, vendor_scorecard, dora_register, benchmark.
Files are saved to MinIO object storage.
"""

import io
import logging
import time
from typing import Any

from app.agents.base_agent import AgentResult, BaseAgent
from app.agents.celery_app import celery_app
from app.config import settings

logger = logging.getLogger("cyberscore.agents.report")

REPORT_TYPES = {
    "executive": {
        "name": "Rapport Exécutif COMEX",
        "default_format": "pptx",
        "template": "executive_report.html",
    },
    "rssi": {
        "name": "Rapport RSSI Détaillé",
        "default_format": "pdf",
        "template": "rssi_report.html",
    },
    "vendor_scorecard": {
        "name": "Scorecard Fournisseur",
        "default_format": "pdf",
        "template": "vendor_scorecard.html",
    },
    "dora_register": {
        "name": "Registre DORA (ACPR)",
        "default_format": "xlsx",
        "template": "dora_register.html",
    },
    "benchmark": {
        "name": "Rapport Benchmark Sectoriel",
        "default_format": "pdf",
        "template": "rssi_report.html",
    },
}


class ReportAgent(BaseAgent):
    """Generates branded reports in PDF, PPTX, and Excel formats."""

    def __init__(self) -> None:
        super().__init__(name="report", timeout=120.0)

    async def execute(self, vendor_id: str, **kwargs: Any) -> AgentResult:
        """Generate a report.

        Args:
            vendor_id: Vendor UUID (nullable for portfolio reports).
            **kwargs: report_type, format, user_id.

        Returns:
            AgentResult with file path and metadata.
        """
        report_type = kwargs.get("report_type", "vendor_scorecard")
        fmt = kwargs.get("format", REPORT_TYPES[report_type]["default_format"])
        user_id = kwargs.get("user_id", "system")
        start = time.monotonic()

        self.logger.info(
            "Generating %s report (%s) for vendor %s",
            report_type,
            fmt,
            vendor_id,
        )

        try:
            if fmt == "pdf":
                file_path = await self._generate_pdf(report_type, vendor_id)
            elif fmt == "pptx":
                file_path = await self._generate_pptx(report_type, vendor_id)
            elif fmt == "xlsx":
                file_path = await self._generate_xlsx(report_type, vendor_id)
            else:
                file_path = await self._generate_pdf(report_type, vendor_id)

            duration = time.monotonic() - start
            return AgentResult(
                agent_name=self.name,
                vendor_id=vendor_id,
                success=True,
                data={
                    "file_path": file_path,
                    "report_type": report_type,
                    "format": fmt,
                    "generated_by": user_id,
                },
                duration_seconds=round(duration, 2),
            )
        except Exception as exc:
            duration = time.monotonic() - start
            return AgentResult(
                agent_name=self.name,
                vendor_id=vendor_id,
                success=False,
                errors=[str(exc)],
                duration_seconds=round(duration, 2),
            )

    def _get_template_env(self) -> "Environment":
        """Create Jinja2 environment loading from templates/reports/."""
        import os

        from jinja2 import Environment, FileSystemLoader, select_autoescape

        template_dir = os.path.join(
            os.path.dirname(os.path.dirname(__file__)),
            "templates",
            "reports",
        )
        return Environment(
            loader=FileSystemLoader(template_dir),
            autoescape=select_autoescape(["html"]),
        )

    def _build_template_context(
        self, report_type: str, vendor_id: str,
    ) -> dict[str, Any]:
        """Build context variables for a given report type."""
        generated_at = time.strftime("%Y-%m-%d %H:%M:%S")
        base = {
            "vendor_id": vendor_id,
            "report_type": report_type,
            "report_name": REPORT_TYPES[report_type]["name"],
            "generated_at": generated_at,
        }

        if report_type == "executive":
            base.update(
                score=580,
                grade="C",
                grade_color="#F39C12",
                date=time.strftime("%d/%m/%Y"),
                dora_coverage=72,
                trends={"delta": 15, "improved": 8, "degraded": 3, "tier1_count": 12},
                top_risks=[
                    {"vendor_name": "CloudCorp", "domain": "D7", "severity": "critical", "description": "Fuites de donnees detectees sur le dark web"},
                    {"vendor_name": "DataHost", "domain": "D3", "severity": "high", "description": "Certificat TLS expire depuis 15 jours"},
                    {"vendor_name": "NetServ", "domain": "D1", "severity": "high", "description": "Ports critiques exposes (22, 3389)"},
                ],
                top_actions=[
                    {"severity": "critical", "title": "Remediation fuites CloudCorp", "description": "Contacter le fournisseur pour plan de remediation sous 48h"},
                    {"severity": "high", "title": "Renouvellement TLS DataHost", "description": "Exiger renouvellement certificat et mise en place monitoring"},
                    {"severity": "high", "title": "Durcissement reseau NetServ", "description": "Fermeture ports non essentiels et mise en place filtrage"},
                ],
            )
        elif report_type == "rssi":
            base.update(
                vendor_name=vendor_id or "Portfolio",
                domain_scores=[
                    {"code": "D1", "name": "Securite Reseau", "score": 62, "grade": "B"},
                    {"code": "D2", "name": "Securite DNS", "score": 72, "grade": "B"},
                    {"code": "D3", "name": "Securite Web", "score": 55, "grade": "C"},
                    {"code": "D4", "name": "Securite Email", "score": 68, "grade": "B"},
                    {"code": "D5", "name": "Cadence Correctifs", "score": 48, "grade": "C"},
                    {"code": "D6", "name": "Reputation IP", "score": 74, "grade": "B"},
                    {"code": "D7", "name": "Fuites & Exposition", "score": 42, "grade": "C"},
                    {"code": "D8", "name": "Presence Reglementaire", "score": 78, "grade": "B"},
                ],
                findings=[
                    {"domain": "D1", "severity": "high", "title": "Ports SSH exposes", "source": "Shodan", "impact": "Acces distant non securise"},
                    {"domain": "D3", "severity": "critical", "title": "Certificat TLS expire", "source": "Scan SSL", "impact": "Donnees en transit non protegees"},
                    {"domain": "D7", "severity": "critical", "title": "Credentials fuites", "source": "HIBP", "impact": "Compromission potentielle de comptes"},
                ],
                recommendations=[
                    {"domain": "D1", "priority": "haute", "text": "Fermer les ports SSH publics et utiliser un VPN pour l'acces distant"},
                    {"domain": "D3", "priority": "critique", "text": "Renouveler immediatement le certificat TLS et activer le renouvellement automatique"},
                    {"domain": "D7", "priority": "critique", "text": "Forcer la rotation de tous les mots de passe compromis"},
                ],
            )
        elif report_type == "vendor_scorecard":
            base.update(
                vendor_name=vendor_id or "Fournisseur",
                score=640,
                grade="B",
                grade_color="#2ECC71",
                tier="Tier 1",
                sector="Assurance",
                last_scan=time.strftime("%d/%m/%Y"),
                domain_scores=[
                    {"code": "D1", "name": "Securite Reseau", "score": 70, "grade": "B"},
                    {"code": "D2", "name": "Securite DNS", "score": 72, "grade": "B"},
                    {"code": "D3", "name": "Securite Web", "score": 65, "grade": "B"},
                    {"code": "D4", "name": "Securite Email", "score": 68, "grade": "B"},
                    {"code": "D5", "name": "Cadence Correctifs", "score": 58, "grade": "C"},
                    {"code": "D6", "name": "Reputation IP", "score": 74, "grade": "B"},
                    {"code": "D7", "name": "Fuites & Exposition", "score": 55, "grade": "C"},
                    {"code": "D8", "name": "Presence Reglementaire", "score": 76, "grade": "B"},
                ],
                trend_delta=12,
                trend_direction="trend-up",
                key_findings=[
                    {"severity": "high", "title": "Ports exposes", "description": "3 ports critiques accessibles publiquement"},
                    {"severity": "medium", "title": "SPF incomplet", "description": "Politique SPF trop permissive (~all)"},
                ],
                recommendations=[
                    "Fermer les ports non essentiels et mettre en place un filtrage reseau",
                    "Durcir la politique SPF avec -all et activer DMARC reject",
                    "Mettre en place un programme de gestion des correctifs avec SLA 30 jours",
                ],
            )
        elif report_type == "dora_register":
            base.update(
                register_ref=f"DORA-REG-{time.strftime('%Y%m%d')}",
                summary={
                    "total_providers": 24,
                    "critical_count": 8,
                    "with_subcontractors": 15,
                    "with_exit_plans": 18,
                },
                providers=[
                    {
                        "name": "CloudCorp SAS", "ict_function": "Hebergement Cloud IaaS",
                        "criticality": "critical", "contract_start": "2022-01-15",
                        "contract_end": "2025-01-14", "status": "active",
                        "subcontractors": ["AWS EU", "OVH Backup"],
                        "certifications": ["ISO 27001", "SOC 2", "HDS"],
                        "data_location": "France / UE", "exit_plan": True,
                        "last_audit": "2024-06-15", "score": 720,
                    },
                    {
                        "name": "DataHost SARL", "ict_function": "Hebergement Donnees Sante",
                        "criticality": "critical", "contract_start": "2023-03-01",
                        "contract_end": "2026-02-28", "status": "active",
                        "subcontractors": [],
                        "certifications": ["ISO 27001", "HDS", "SecNumCloud"],
                        "data_location": "France", "exit_plan": True,
                        "last_audit": "2024-09-20", "score": 810,
                    },
                    {
                        "name": "NetServ SA", "ict_function": "Reseau & Telecommunications",
                        "criticality": "important", "contract_start": "2021-06-01",
                        "contract_end": "2024-05-31", "status": "review",
                        "subcontractors": ["Orange Business"],
                        "certifications": ["ISO 27001"],
                        "data_location": "France", "exit_plan": False,
                        "last_audit": "2023-12-10", "score": 540,
                    },
                ],
            )

        return base

    async def _generate_pdf(
        self, report_type: str, vendor_id: str
    ) -> str:
        """Generate PDF via WeasyPrint (HTML template -> PDF).

        Renders a Jinja2 HTML template with vendor data and converts
        to PDF via WeasyPrint. Uploads the result to MinIO.
        """
        from weasyprint import HTML

        env = self._get_template_env()
        template_name = REPORT_TYPES[report_type]["template"]
        template = env.get_template(template_name)
        context = self._build_template_context(report_type, vendor_id)
        html_content = template.render(**context)

        pdf_bytes = HTML(string=html_content).write_pdf()
        object_key = f"reports/{vendor_id}/{report_type}.pdf"
        await self._upload_to_minio(object_key, pdf_bytes, "application/pdf")
        return object_key

    async def _generate_pptx(
        self, report_type: str, vendor_id: str
    ) -> str:
        """Generate branded PPTX via python-pptx."""
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Emu, Inches, Pt

        NAVY = RGBColor(0x1B, 0x3A, 0x5C)
        BLUE = RGBColor(0x2E, 0x75, 0xB6)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        TEXT_COLOR = RGBColor(0x2C, 0x3E, 0x50)
        LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFA)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        context = self._build_template_context(report_type, vendor_id)

        def add_title_bar(slide, height_inches=0.8):
            """Add a navy title bar at the top of a slide."""
            from pptx.util import Inches as In
            shape = slide.shapes.add_shape(
                1, In(0), In(0), prs.slide_width, In(height_inches),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = NAVY
            shape.line.fill.background()
            return shape

        # --- Slide 1: Title ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        add_title_bar(slide, 1.2)

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(10), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "CyberScore"
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = WHITE

        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(10), Inches(2))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = REPORT_TYPES[report_type]["name"]
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = NAVY

        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"Fournisseur: {vendor_id}\nDate: {context['generated_at']}"
        run.font.size = Pt(18)
        run.font.color.rgb = TEXT_COLOR

        # --- Slide 2: Score Overview ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bar = add_title_bar(slide)

        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8), Inches(0.5))
        tf = header_box.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = "Score Global"
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = WHITE

        score_box = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(5), Inches(3))
        tf = score_box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = str(context.get("score", "—"))
        run.font.size = Pt(72)
        run.font.bold = True
        run.font.color.rgb = BLUE

        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"Grade: {context.get('grade', '—')}"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = NAVY

        # --- Slide 3: Domain scores ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title_bar(slide)

        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8), Inches(0.5))
        tf = header_box.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = "Scores par Domaine"
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = WHITE

        domain_scores = context.get("domain_scores", [])
        for i, ds in enumerate(domain_scores[:8]):
            y = 1.2 + i * 0.7
            label_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(3), Inches(0.5))
            tf = label_box.text_frame
            run = tf.paragraphs[0].add_run()
            name = ds.get("name", ds.get("code", ""))
            run.text = f"{ds.get('code', '')} {name}"
            run.font.size = Pt(14)
            run.font.color.rgb = TEXT_COLOR

            score_label = slide.shapes.add_textbox(Inches(10), Inches(y), Inches(2), Inches(0.5))
            tf = score_label.text_frame
            run = tf.paragraphs[0].add_run()
            run.text = f"{ds.get('score', 0)}/100 ({ds.get('grade', '-')})"
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = BLUE

        buffer = io.BytesIO()
        prs.save(buffer)
        pptx_bytes = buffer.getvalue()

        object_key = f"reports/{vendor_id}/{report_type}.pptx"
        await self._upload_to_minio(
            object_key, pptx_bytes,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        return object_key

    async def _generate_xlsx(
        self, report_type: str, vendor_id: str
    ) -> str:
        """Generate Excel via openpyxl."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = REPORT_TYPES[report_type]["name"]

        # Header row
        headers = ["Fournisseur", "Domaine", "Score", "Grade", "Date"]
        ws.append(headers)

        # Placeholder data row
        ws.append([vendor_id, "-", "-", "-", time.strftime("%Y-%m-%d")])

        buffer = io.BytesIO()
        wb.save(buffer)
        xlsx_bytes = buffer.getvalue()

        object_key = f"reports/{vendor_id}/{report_type}.xlsx"
        await self._upload_to_minio(
            object_key, xlsx_bytes,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
        return object_key

    async def _upload_to_minio(
        self, object_key: str, data: bytes, content_type: str,
    ) -> None:
        """Upload a file to MinIO object storage.

        Args:
            object_key: S3-style object key (path).
            data: File bytes.
            content_type: MIME type.
        """
        try:
            from minio import Minio

            client = Minio(
                settings.minio_endpoint,
                access_key=settings.minio_access_key,
                secret_key=settings.minio_secret_key,
                secure=False,
            )
            bucket = settings.minio_bucket
            if not client.bucket_exists(bucket):
                client.make_bucket(bucket)

            client.put_object(
                bucket,
                object_key,
                io.BytesIO(data),
                length=len(data),
                content_type=content_type,
            )
            self.logger.info("Uploaded %s to MinIO (%d bytes)", object_key, len(data))
        except Exception as exc:
            self.logger.warning(
                "MinIO upload failed for %s, saving locally: %s", object_key, exc
            )
            # Fallback: save locally
            import os
            local_path = os.path.join("data", object_key)
            os.makedirs(os.path.dirname(local_path), exist_ok=True)
            with open(local_path, "wb") as f:
                f.write(data)


@celery_app.task(name="app.agents.report_agent.generate_report")
def generate_report(
    report_type: str,
    vendor_id: str = "",
    user_id: str = "system",
    fmt: str = "",
) -> dict[str, Any]:
    """Celery task: generate a report."""
    import asyncio

    agent = ReportAgent()
    result = asyncio.run(
        agent.execute(
            vendor_id,
            report_type=report_type,
            format=fmt or REPORT_TYPES[report_type]["default_format"],
            user_id=user_id,
        )
    )
    return result.data
